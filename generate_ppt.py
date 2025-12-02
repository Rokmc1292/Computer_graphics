#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
3D 군 생활관 시뮬레이션 프로젝트 PPT 자동 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from io import BytesIO

def create_scene_graph():
    """Scene Graph 시각화 이미지 생성"""
    fig, ax = plt.subplots(figsize=(14, 10))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 10)
    ax.axis('off')

    # 색상 정의
    color_scene = '#E8F4F8'
    color_manager = '#B8E6F5'
    color_object = '#FFE5B4'
    color_control = '#D4F1D4'
    color_loader = '#F5D4F5'

    # 노드 정의 (x, y, width, height, text, color)
    nodes = [
        # Root
        (6, 9, 2, 0.6, 'Scene', color_scene),

        # Managers
        (2, 7.5, 2, 0.5, 'SceneManager', color_manager),
        (5, 7.5, 2, 0.5, 'LightManager', color_manager),
        (8, 7.5, 2, 0.5, 'InteractionController', color_control),

        # Objects Group
        (1, 6, 2, 0.5, 'Barracks', color_object),
        (3.5, 6, 1.5, 0.5, 'Avatar', color_object),

        # Barracks Children - Row 1
        (0.5, 4.5, 1.3, 0.4, 'Floor', color_object),
        (2, 4.5, 1.3, 0.4, 'Walls', color_object),
        (3.5, 4.5, 1.3, 0.4, 'Ceiling', color_object),
        (5, 4.5, 1.3, 0.4, 'Particles', color_object),

        # Barracks Children - Row 2
        (0.2, 3, 1.1, 0.4, 'BunkBed (x6)', color_object),
        (1.5, 3, 1.1, 0.4, 'Chester (x8)', color_object),
        (2.8, 3, 1.1, 0.4, 'Locker', color_object),
        (4.1, 3, 1.1, 0.4, 'Window (x4)', color_object),
        (5.4, 3, 1.1, 0.4, 'Door', color_object),

        # Barracks Children - Row 3
        (0.5, 1.5, 1.1, 0.4, 'TV', color_object),
        (1.8, 1.5, 1.3, 0.4, 'Radiator (x2)', color_object),
        (3.3, 1.5, 1.3, 0.4, 'CeilingFan', color_object),
        (4.8, 1.5, 1.3, 0.4, 'LightSwitch', color_object),

        # Barracks Children - Row 4
        (0.3, 0.2, 1.5, 0.4, 'DecorativeCurve', color_object),
        (2, 0.2, 1.2, 0.4, 'Vase (x6)', color_object),

        # Controls
        (9, 6, 1.8, 0.5, 'AvatarController', color_control),
        (11, 6, 1.8, 0.5, 'CameraController', color_control),

        # Loaders
        (9, 4.5, 1.8, 0.5, 'TextureLoader', color_loader),
        (11, 4.5, 1.8, 0.5, 'ModelLoader', color_loader),
    ]

    # 노드 그리기
    for x, y, w, h, text, color in nodes:
        rect = patches.FancyBboxPatch((x, y), w, h,
                                       boxstyle="round,pad=0.05",
                                       edgecolor='black',
                                       facecolor=color,
                                       linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h/2, text,
                ha='center', va='center',
                fontsize=8, fontweight='bold')

    # 연결선 그리기
    connections = [
        # Scene to Managers
        (7, 9, 3, 8),
        (7, 9, 6, 8),
        (7, 9, 9, 8),

        # Barracks to children
        (2, 6, 1.15, 4.7),
        (2, 6, 2.65, 4.7),
        (2, 6, 4.15, 4.7),
        (2, 6, 5.65, 4.7),

        (2, 6, 0.75, 3.4),
        (2, 6, 2.05, 3.4),
        (2, 6, 3.35, 3.4),
        (2, 6, 4.65, 3.4),
        (2, 6, 5.95, 3.4),

        (2, 6, 1.05, 1.9),
        (2, 6, 2.45, 1.9),
        (2, 6, 3.95, 1.9),
        (2, 6, 5.45, 1.9),

        (2, 6, 1.05, 0.6),
        (2, 6, 2.6, 0.6),

        # Avatar
        (7, 9, 4.25, 6.5),

        # Controls
        (7, 9, 9.9, 6.5),
        (7, 9, 11.9, 6.5),

        # Loaders
        (6, 7.5, 9.9, 5),
        (6, 7.5, 11.9, 5),
    ]

    for x1, y1, x2, y2 in connections:
        ax.plot([x1, x2], [y1, y2], 'k-', linewidth=1, alpha=0.5)

    # 범례
    legend_elements = [
        patches.Patch(facecolor=color_scene, edgecolor='black', label='Scene Root'),
        patches.Patch(facecolor=color_manager, edgecolor='black', label='Scene/Light Managers'),
        patches.Patch(facecolor=color_object, edgecolor='black', label='3D Objects'),
        patches.Patch(facecolor=color_control, edgecolor='black', label='Controllers'),
        patches.Patch(facecolor=color_loader, edgecolor='black', label='Loaders')
    ]
    ax.legend(handles=legend_elements, loc='upper right', fontsize=9)

    plt.title('3D 군생활관 Scene Graph', fontsize=16, fontweight='bold', pad=10)
    plt.tight_layout()

    # 이미지를 BytesIO로 저장
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
    img_stream.seek(0)
    plt.close()

    return img_stream

def create_title_slide(prs):
    """타이틀 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "3D 군 생활관 시뮬레이션"
    subtitle.text = "Three.js 기반 인터랙티브 3D 웹 애플리케이션"

    # 스타일 적용
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)

def create_overview_slide(prs):
    """프로젝트 개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "프로젝트 개요"

    body = slide.placeholders[1].text_frame
    body.text = "목적"
    p = body.add_paragraph()
    p.text = "웹 브라우저에서 실행되는 실감나는 3D 군 생활관 환경 구현"
    p.level = 1

    p = body.add_paragraph()
    p.text = "주요 특징"

    p = body.add_paragraph()
    p.text = "인터랙티브한 아바타 조작 (WASD 이동)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "PBR(Physically Based Rendering) 기반 사실적인 렌더링"
    p.level = 1

    p = body.add_paragraph()
    p.text = "실시간 애니메이션 (조명, 창문, 문 등)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "고품질 그림자 및 조명 시스템"
    p.level = 1

def create_tech_stack_slide(prs):
    """기술 스택"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "기술 스택"

    body = slide.placeholders[1].text_frame
    body.text = "핵심 프레임워크"

    p = body.add_paragraph()
    p.text = "Three.js v0.161.0 - 3D 그래픽 라이브러리"
    p.level = 1

    p = body.add_paragraph()
    p.text = "WebGL Renderer - 하드웨어 가속 렌더링"
    p.level = 1

    p = body.add_paragraph()
    p.text = "렌더링 기술"

    p = body.add_paragraph()
    p.text = "PBR Materials (MeshStandardMaterial)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "ACES Filmic Tone Mapping"
    p.level = 1

    p = body.add_paragraph()
    p.text = "PCF Soft Shadows (4096x4096 해상도)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "텍스처 맵: Diffuse, Normal, Roughness"
    p.level = 1

def create_architecture_slide(prs):
    """프로젝트 아키텍처"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "프로젝트 아키텍처"

    body = slide.placeholders[1].text_frame
    body.text = "모듈 구조"

    modules = [
        ("main.js", "애플리케이션 부트스트랩 및 초기화"),
        ("scene/SceneManager.js", "Three.js 씬, 카메라, 렌더러 설정"),
        ("scene/LightManager.js", "조명 시스템 관리"),
        ("objects/Barracks.js", "생활관 메인 오케스트레이터"),
        ("objects/Avatar.js", "플레이어 캐릭터"),
        ("controls/AvatarController.js", "WASD 키보드 입력 처리"),
        ("controls/CameraController.js", "마우스 카메라 제어"),
        ("controls/InteractionController.js", "객체 상호작용 관리"),
        ("loaders/", "텍스처 및 3D 모델 로딩")
    ]

    for module, desc in modules:
        p = body.add_paragraph()
        p.text = f"{module}: {desc}"
        p.level = 1
        p.font.size = Pt(15)

def create_scene_graph_slide(prs):
    """Scene Graph 시각화 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    # 타이틀 추가
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Scene Graph 구조"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(31, 73, 125)
    title_para.alignment = PP_ALIGN.CENTER

    # Scene Graph 이미지 생성 및 추가
    img_stream = create_scene_graph()
    pic = slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1), width=Inches(9))

    print("✓ Scene Graph 슬라이드 생성 완료")

def create_scene_management_slide(prs):
    """Scene Management"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Scene Management"

    body = slide.placeholders[1].text_frame
    body.text = "SceneManager (scene/SceneManager.js)"

    p = body.add_paragraph()
    p.text = "Three.js Scene 초기화 및 설정"
    p.level = 1

    p = body.add_paragraph()
    p.text = "PerspectiveCamera 구성 (FOV: 75°)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "WebGL Renderer 설정 (안티앨리어싱, 그림자)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "반응형 윈도우 리사이즈 처리"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "LightManager (scene/LightManager.js)"

    p = body.add_paragraph()
    p.text = "Ambient Light - 부드러운 간접 조명"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Directional Light - 태양광 (창문을 통한 빛)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Point Lights - 천장 조명 4개"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Spotlights - 창문 효과 4개"
    p.level = 1

def create_objects_slide1(prs):
    """객체 시스템 - 방 구조"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "객체 시스템: 생활관 구조"

    body = slide.placeholders[1].text_frame
    body.text = "Barracks.js - 메인 오케스트레이터"

    p = body.add_paragraph()
    p.text = "방 구조"

    p = body.add_paragraph()
    p.text = "바닥: 대리석 텍스처 (PBR 맵핑)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "벽 4개: 뒷벽에 문 구멍 (3m x 5m)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "천장: 절차적 캔버스 텍스처"
    p.level = 1

    p = body.add_paragraph()
    p.text = "먼지 파티클 시스템 (200개)"
    p.level = 1

def create_objects_slide2(prs):
    """객체 시스템 - 가구"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "객체 시스템: 가구 및 오브젝트"

    body = slide.placeholders[1].text_frame
    body.text = "주요 가구 (총 30개 이상 오브젝트)"

    furniture = [
        ("이층 침대 (BunkBed.js)", "6개 - 침대 시트 바람 애니메이션"),
        ("서랍장 (Chester.js)", "8개 - 각 침대 옆, 서랍 애니메이션"),
        ("사물함 (Locker.js)", "1개 - 개인 물품 보관"),
        ("창문 (Window.js)", "4개 - 주/야간 색상 애니메이션"),
        ("문 (Door.js)", "1개 - 근접 시 자동 개폐"),
        ("TV (TV.js)", "1개 - 화면 깜빡임 효과"),
        ("라디에이터 (Radiator.js)", "2개 - 발열 효과"),
        ("천장 선풍기 (CeilingFan.js)", "1개 - 회전 애니메이션"),
        ("전등 스위치 (LightSwitch.js)", "1개 - 클릭 상호작용"),
        ("장식 커브 (DecorativeCurve.js)", "금색 Bezier 커브 장식"),
        ("화병 (Vase.js)", "6개 - LatheGeometry 회전체")
    ]

    for name, desc in furniture:
        p = body.add_paragraph()
        p.text = f"{name}: {desc}"
        p.level = 1
        p.font.size = Pt(14)

def create_curve_surface_slide(prs):
    """Curve & Surface 지오메트리"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "고급 지오메트리: Curve & Surface"

    body = slide.placeholders[1].text_frame
    body.text = "DecorativeCurve - Bezier Curve (objects/DecorativeCurve.js)"

    p = body.add_paragraph()
    p.text = "CubicBezierCurve3로 3D 곡선 생성"
    p.level = 1

    p = body.add_paragraph()
    p.text = "TubeGeometry로 곡선을 따라 튜브 형성"
    p.level = 1

    p = body.add_paragraph()
    p.text = "발광 효과가 있는 금색 장식 트림 (앞벽)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "4개 제어점으로 부드러운 곡선 표현"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "Vase - LatheGeometry (objects/Vase.js)"

    p = body.add_paragraph()
    p.text = "LatheGeometry로 회전체 생성 (화병, 컵)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "2D 프로파일을 회전축 중심으로 회전시켜 3D 형상"
    p.level = 1

    p = body.add_paragraph()
    p.text = "서랍장 위에 6개 배치 (다양한 크기)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "도자기 재질 시뮬레이션 (높은 roughness)"
    p.level = 1

def create_avatar_controls_slide(prs):
    """아바타 & 컨트롤"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "아바타 & 컨트롤 시스템"

    body = slide.placeholders[1].text_frame
    body.text = "Avatar (objects/Avatar.js)"

    p = body.add_paragraph()
    p.text = "머리, 몸통, 팔, 다리로 구성된 심플한 캐릭터"
    p.level = 1

    p = body.add_paragraph()
    p.text = "군복 색상 적용 (밀리터리 그린: #4A5D23)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "그림자 캐스팅 지원"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "AvatarController (controls/AvatarController.js)"

    p = body.add_paragraph()
    p.text = "WASD 키보드 입력으로 이동"
    p.level = 1

    p = body.add_paragraph()
    p.text = "카메라 상대 방향 이동"
    p.level = 1

    p = body.add_paragraph()
    p.text = "방 경계 충돌 감지 (±9, ±7 units)"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "CameraController (controls/CameraController.js)"

    p = body.add_paragraph()
    p.text = "마우스 드래그로 카메라 회전 (요/피치)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "마우스 휠로 줌 (3-15 units)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "3인칭 뷰로 아바타 추적"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "InteractionController (controls/InteractionController.js)"

    p = body.add_paragraph()
    p.text = "마우스 클릭으로 객체와 상호작용"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Raycasting으로 클릭 가능한 객체 감지"
    p.level = 1

    p = body.add_paragraph()
    p.text = "전등 스위치 토글 기능"
    p.level = 1

def create_rendering_slide(prs):
    """렌더링 기술"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "렌더링 기술"

    body = slide.placeholders[1].text_frame
    body.text = "PBR (Physically Based Rendering)"

    p = body.add_paragraph()
    p.text = "MeshStandardMaterial 사용"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Diffuse/Albedo 맵 - 기본 색상"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Normal 맵 - 표면 디테일"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Roughness 맵 - 거칠기"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Metalness 파라미터"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "고급 렌더링 설정"

    p = body.add_paragraph()
    p.text = "ACES Filmic Tone Mapping - 영화 같은 색감"
    p.level = 1

    p = body.add_paragraph()
    p.text = "sRGB 색공간 인코딩"
    p.level = 1

    p = body.add_paragraph()
    p.text = "PCF Soft Shadows - 부드러운 그림자"
    p.level = 1

    p = body.add_paragraph()
    p.text = "4096x4096 그림자 맵 해상도"
    p.level = 1

    p = body.add_paragraph()
    p.text = "안티앨리어싱 및 고해상도 픽셀 비율 지원"
    p.level = 1

def create_animation_slide(prs):
    """애니메이션 시스템"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "애니메이션 시스템"

    body = slide.placeholders[1].text_frame
    body.text = "실시간 애니메이션 (총 8개)"

    animations = [
        ("천장 조명 깜빡임", "빠른 플리커 + 느린 펄스 + 랜덤 노이즈"),
        ("TV 화면 펄싱", "청색 발광 효과 애니메이션"),
        ("창문 색상 사이클", "주간 → 저녁 → 야간 색상 변화"),
        ("침대 시트 바람", "UV 오프셋 애니메이션"),
        ("천장 선풍기 회전", "연속 회전 애니메이션"),
        ("문 개폐", "근접 시 부드러운 회전 (< 3 units)"),
        ("먼지 파티클", "브라운 운동 시뮬레이션"),
        ("라디에이터 발열", "열 효과 애니메이션")
    ]

    for name, desc in animations:
        p = body.add_paragraph()
        p.text = f"{name}: {desc}"
        p.level = 1
        p.font.size = Pt(16)

def create_features_slide(prs):
    """주요 기능"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "주요 기능 요약"

    body = slide.placeholders[1].text_frame
    body.text = "인터랙션"

    p = body.add_paragraph()
    p.text = "WASD로 아바타 이동 (카메라 상대 방향)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "마우스 드래그/휠로 카메라 조작"
    p.level = 1

    p = body.add_paragraph()
    p.text = "자동 문 개폐 (근접 감지)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "전등 스위치 클릭 상호작용"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "시각 효과"

    p = body.add_paragraph()
    p.text = "PBR 재질과 텍스처 맵핑"
    p.level = 1

    p = body.add_paragraph()
    p.text = "9개 광원 (포인트, 스팟, 디렉셔널)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "8개 실시간 애니메이션"
    p.level = 1

    p = body.add_paragraph()
    p.text = "고품질 그림자 (4K 해상도)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Bezier Curve 및 LatheGeometry 장식"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "자산 관리"

    p = body.add_paragraph()
    p.text = "GLB 3D 모델 로딩 (5개)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "PBR 텍스처 세트 (대리석, 리넨)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "에러 핸들링 및 타임아웃 처리"
    p.level = 1

def create_code_stats_slide(prs):
    """코드 통계"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "프로젝트 통계"

    body = slide.placeholders[1].text_frame
    body.text = "코드 규모"

    p = body.add_paragraph()
    p.text = "총 21개 JavaScript 파일"
    p.level = 1

    p = body.add_paragraph()
    p.text = "약 2,000줄 이상의 코드"
    p.level = 1

    p = body.add_paragraph()
    p.text = "프로젝트 크기: ~80MB"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "3D 자산"

    p = body.add_paragraph()
    p.text = "5개 GLB 3D 모델 (문, 창문, TV, 선풍기, 라디에이터)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "2개 PBR 텍스처 세트 (대리석, 리넨)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "총 30개 이상 씬 오브젝트"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "모듈 구성"

    p = body.add_paragraph()
    p.text = "Scene: 2개 (SceneManager, LightManager)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Controls: 3개 (Avatar, Camera, Interaction)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Objects: 12개 (Barracks, Avatar, 가구, Curve, Vase 등)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "Loaders: 2개 (Texture, Model)"
    p.level = 1

def create_conclusion_slide(prs):
    """결론"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "프로젝트 특징 및 성과"

    body = slide.placeholders[1].text_frame
    body.text = "기술적 성과"

    p = body.add_paragraph()
    p.text = "고급 Three.js 기능 활용 (조명, 재질, 그림자, 애니메이션)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "모듈형 아키텍처로 명확한 관심사 분리"
    p.level = 1

    p = body.add_paragraph()
    p.text = "자산 로딩 및 에러 핸들링 구현"
    p.level = 1

    p = body.add_paragraph()
    p.text = "실시간 애니메이션 시스템 구축"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "사용자 경험"

    p = body.add_paragraph()
    p.text = "직관적인 WASD + 마우스 컨트롤"
    p.level = 1

    p = body.add_paragraph()
    p.text = "몰입감 있는 3D 환경"
    p.level = 1

    p = body.add_paragraph()
    p.text = "실시간 인터랙션 (문 개폐 등)"
    p.level = 1

    p = body.add_paragraph()
    p.text = ""

    p = body.add_paragraph()
    p.text = "향후 개선 방향"

    p = body.add_paragraph()
    p.text = "추가 인터랙션 요소 (서랍, 사물함 열기)"
    p.level = 1

    p = body.add_paragraph()
    p.text = "더 많은 애니메이션 및 효과"
    p.level = 1

    p = body.add_paragraph()
    p.text = "성능 최적화 (LOD, Instancing)"
    p.level = 1

def main():
    """메인 함수"""
    print("🎨 PPT 생성 시작...")

    # 프레젠테이션 객체 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 슬라이드 생성
    print("📄 타이틀 슬라이드 생성...")
    create_title_slide(prs)

    print("📄 프로젝트 개요 슬라이드 생성...")
    create_overview_slide(prs)

    print("📄 기술 스택 슬라이드 생성...")
    create_tech_stack_slide(prs)

    print("📄 프로젝트 아키텍처 슬라이드 생성...")
    create_architecture_slide(prs)

    print("📄 Scene Graph 슬라이드 생성...")
    create_scene_graph_slide(prs)

    print("📄 Scene Management 슬라이드 생성...")
    create_scene_management_slide(prs)

    print("📄 객체 시스템 슬라이드 생성...")
    create_objects_slide1(prs)
    create_objects_slide2(prs)

    print("📄 Curve & Surface 슬라이드 생성...")
    create_curve_surface_slide(prs)

    print("📄 아바타 & 컨트롤 슬라이드 생성...")
    create_avatar_controls_slide(prs)

    print("📄 렌더링 기술 슬라이드 생성...")
    create_rendering_slide(prs)

    print("📄 애니메이션 시스템 슬라이드 생성...")
    create_animation_slide(prs)

    print("📄 주요 기능 슬라이드 생성...")
    create_features_slide(prs)

    print("📄 프로젝트 통계 슬라이드 생성...")
    create_code_stats_slide(prs)

    print("📄 결론 슬라이드 생성...")
    create_conclusion_slide(prs)

    # 파일 저장
    output_file = '3D_군생활관_프로젝트_설명.pptx'
    prs.save(output_file)

    print(f"✅ PPT 생성 완료: {output_file}")
    print(f"   총 {len(prs.slides)} 슬라이드 생성됨")

if __name__ == "__main__":
    main()
